import io

from pptx import Presentation
from pptx.util import Inches

from core.ingest import _build_document_from_pptx_fallback


def test_pptx_fallback_extracts_text_blocks_and_renders_a_page() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    text_box.text_frame.text = "Attention maps a query to relevant keys."
    file_bytes = io.BytesIO()
    presentation.save(file_bytes)

    document = _build_document_from_pptx_fallback(file_bytes.getvalue(), "demo.pptx", 110)

    assert len(document.pages) == 1
    assert document.pages[0].blocks[0].text == "Attention maps a query to relevant keys."
    assert document.pages[0].image_png.startswith(b"\x89PNG")
