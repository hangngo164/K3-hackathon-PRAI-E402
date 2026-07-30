"""PDF and PowerPoint ingestion for CP2.

This module intentionally has no Streamlit dependency so it can later be
reused by the evaluation runner.
"""

from __future__ import annotations

import hashlib
import io
import shutil
import subprocess
import tempfile
import textwrap
import winreg
from pathlib import Path

import fitz
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from core.models import Block, Document, Page


EMU_PER_INCH = 914400
POINTS_PER_INCH = 72


class OfficeConversionUnavailable(RuntimeError):
    """Raised when no local PowerPoint conversion backend is available."""


def build_document(file_bytes: bytes, source_name: str, dpi: int) -> Document:
    """Prepare PDF, PPTX, or legacy PPT input for the selectable-block viewer."""

    suffix = Path(source_name).suffix.lower()
    if suffix == ".pdf":
        return _build_document_from_pdf(file_bytes, source_name, dpi)

    if suffix not in {".ppt", ".pptx"}:
        raise ValueError("Chỉ hỗ trợ file PDF, PPTX hoặc PPT.")

    try:
        converted_pdf = _convert_powerpoint_to_pdf(file_bytes, source_name)
    except OfficeConversionUnavailable:
        if suffix == ".pptx":
            return _build_document_from_pptx_fallback(file_bytes, source_name, dpi)
        raise ValueError(
            "File PPT cần Microsoft PowerPoint hoặc LibreOffice để chuyển sang PDF. "
            "Cài một trong hai hoặc xuất file này thành PDF rồi thử lại."
        ) from None

    return _build_document_from_pdf(converted_pdf, source_name, dpi)


def _build_document_from_pdf(file_bytes: bytes, source_name: str, dpi: int) -> Document:
    """Extract selectable text blocks and render every page of a PDF."""

    pdf = fitz.open(stream=file_bytes, filetype="pdf")
    pages: list[Page] = []
    scale = dpi / 72

    try:
        for page_index, pdf_page in enumerate(pdf):
            page_no = page_index + 1
            blocks = _extract_blocks(pdf_page, page_no)
            pixmap = pdf_page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
            pages.append(
                Page(
                    page_no=page_no,
                    blocks=tuple(blocks),
                    text="\n".join(block.text for block in blocks),
                    image_png=pixmap.tobytes("png"),
                    width_pt=pdf_page.rect.width,
                    height_pt=pdf_page.rect.height,
                )
            )
    finally:
        pdf.close()

    return Document(
        doc_hash=hashlib.sha256(file_bytes).hexdigest(),
        source_name=source_name,
        pages=tuple(pages),
    )


def _convert_powerpoint_to_pdf(file_bytes: bytes, source_name: str) -> bytes:
    """Prefer Microsoft PowerPoint, then fall back to LibreOffice."""

    try:
        return _convert_with_microsoft_powerpoint(file_bytes, source_name)
    except OfficeConversionUnavailable:
        return _convert_with_libreoffice(file_bytes, source_name)


def _convert_with_microsoft_powerpoint(file_bytes: bytes, source_name: str) -> bytes:
    """Export a presentation to PDF through the locally installed PowerPoint."""

    if not _powerpoint_installed():
        raise OfficeConversionUnavailable

    safe_name = Path(source_name).name
    with tempfile.TemporaryDirectory(prefix="daily-quiz-powerpoint-") as temp_dir:
        input_path = Path(temp_dir, safe_name)
        output_path = Path(temp_dir, "converted.pdf")
        script_path = Path(temp_dir, "export_pdf.ps1")
        input_path.write_bytes(file_bytes)
        script_path.write_text(_POWERPOINT_EXPORT_SCRIPT, encoding="utf-8")
        result = subprocess.run(
            [
                "powershell.exe",
                "-NoProfile",
                "-NonInteractive",
                "-ExecutionPolicy",
                "Bypass",
                "-File",
                str(script_path),
                "-InputPath",
                str(input_path),
                "-OutputPath",
                str(output_path),
            ],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
        if result.returncode != 0 or not output_path.exists():
            raise ValueError("Microsoft PowerPoint không thể chuyển file này sang PDF.")
        return output_path.read_bytes()


def _convert_with_libreoffice(file_bytes: bytes, source_name: str) -> bytes:
    """Use LibreOffice only on machines without Microsoft PowerPoint."""

    soffice = shutil.which("soffice") or shutil.which("soffice.exe")
    if not soffice:
        raise OfficeConversionUnavailable

    safe_name = Path(source_name).name
    with tempfile.TemporaryDirectory(prefix="daily-quiz-ppt-") as temp_dir:
        input_path = Path(temp_dir, safe_name)
        input_path.write_bytes(file_bytes)
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, str(input_path)],
            capture_output=True,
            text=True,
            timeout=90,
            check=False,
        )
        output_path = input_path.with_suffix(".pdf")
        if result.returncode != 0 or not output_path.exists():
            raise ValueError("Không thể chuyển PowerPoint sang PDF bằng LibreOffice.")
        return output_path.read_bytes()


def _powerpoint_installed() -> bool:
    """Check the registered PowerPoint executable without launching Office."""

    key_paths = (
        r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE",
        r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE",
    )
    for key_path in key_paths:
        try:
            with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, key_path) as key:
                executable, _ = winreg.QueryValueEx(key, None)
            if Path(executable).exists():
                return True
        except OSError:
            continue
    return False


_POWERPOINT_EXPORT_SCRIPT = r"""
param(
    [Parameter(Mandatory = $true)][string]$InputPath,
    [Parameter(Mandatory = $true)][string]$OutputPath
)

$application = $null
$presentation = $null
try {
    $application = New-Object -ComObject PowerPoint.Application
    $presentation = $application.Presentations.Open($InputPath, $true, $false, $false)
    # 32 is the Office constant ppSaveAsPDF.
    $presentation.SaveAs($OutputPath, 32)
}
finally {
    if ($presentation -ne $null) { $presentation.Close() }
    if ($application -ne $null) { $application.Quit() }
    [GC]::Collect()
    [GC]::WaitForPendingFinalizers()
}
"""


def _build_document_from_pptx_fallback(file_bytes: bytes, source_name: str, dpi: int) -> Document:
    """Read PPTX text when LibreOffice is unavailable.

    This fallback deliberately favors selectable source text over visual fidelity.
    PowerPoint files converted through LibreOffice always use the PDF path above.
    """

    presentation = Presentation(io.BytesIO(file_bytes))
    width_pt = presentation.slide_width / EMU_PER_INCH * POINTS_PER_INCH
    height_pt = presentation.slide_height / EMU_PER_INCH * POINTS_PER_INCH
    pixel_width = max(1, round(width_pt * dpi / POINTS_PER_INCH))
    pixel_height = max(1, round(height_pt * dpi / POINTS_PER_INCH))
    scale = dpi / POINTS_PER_INCH
    pages: list[Page] = []

    for page_index, slide in enumerate(presentation.slides):
        page_no = page_index + 1
        image = Image.new("RGB", (pixel_width, pixel_height), "white")
        draw = ImageDraw.Draw(image)
        blocks: list[Block] = []

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue

            x0 = shape.left / EMU_PER_INCH * POINTS_PER_INCH
            y0 = shape.top / EMU_PER_INCH * POINTS_PER_INCH
            x1 = x0 + shape.width / EMU_PER_INCH * POINTS_PER_INCH
            y1 = y0 + shape.height / EMU_PER_INCH * POINTS_PER_INCH
            order = len(blocks) + 1
            blocks.append(
                Block(
                    block_id=f"p{page_no:02d}-b{order:02d}",
                    page_no=page_no,
                    order=order,
                    text=text,
                    bbox=(x0, y0, x1, y1),
                    font_size_max=0.0,
                )
            )
            _draw_fallback_text(draw, text, (x0, y0, x1, y1), scale)

        image_bytes = io.BytesIO()
        image.save(image_bytes, format="PNG")
        pages.append(
            Page(
                page_no=page_no,
                blocks=tuple(blocks),
                text="\n".join(block.text for block in blocks),
                image_png=image_bytes.getvalue(),
                width_pt=width_pt,
                height_pt=height_pt,
            )
        )

    return Document(
        doc_hash=hashlib.sha256(file_bytes).hexdigest(),
        source_name=source_name,
        pages=tuple(pages),
    )


def _draw_fallback_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    bbox: tuple[float, float, float, float],
    scale: float,
) -> None:
    """Draw a simple, readable text-only PPTX preview without claiming fidelity."""

    x0, y0, x1, y1 = (round(value * scale) for value in bbox)
    draw.rectangle((x0, y0, x1, y1), outline=(210, 214, 220), width=1)
    font_size = max(14, min(32, round((y1 - y0) * 0.22)))
    try:
        font = ImageFont.truetype("arial.ttf", font_size)
    except OSError:
        font = ImageFont.load_default()
    chars_per_line = max(18, int((x1 - x0) / max(font_size * 0.55, 1)))
    wrapped = "\n".join(textwrap.wrap(text, width=chars_per_line, replace_whitespace=False))
    draw.multiline_text((x0 + 8, y0 + 6), wrapped, fill=(26, 32, 44), font=font, spacing=4)


def _extract_blocks(pdf_page: fitz.Page, page_no: int) -> list[Block]:
    blocks: list[Block] = []
    raw_page = pdf_page.get_text("dict")

    for raw_block in raw_page["blocks"]:
        if raw_block["type"] != 0:
            continue

        text_parts: list[str] = []
        font_sizes: list[float] = []
        for line in raw_block.get("lines", []):
            line_text = "".join(span["text"] for span in line.get("spans", []))
            if line_text.strip():
                text_parts.append(line_text.strip())
            font_sizes.extend(span["size"] for span in line.get("spans", []))

        text = "\n".join(text_parts).strip()
        if not text:
            continue

        order = len(blocks) + 1
        blocks.append(
            Block(
                block_id=f"p{page_no:02d}-b{order:02d}",
                page_no=page_no,
                order=order,
                text=text,
                bbox=tuple(float(value) for value in raw_block["bbox"]),
                font_size_max=max(font_sizes, default=0.0),
            )
        )

    return blocks
